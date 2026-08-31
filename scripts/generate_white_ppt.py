import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_nexora_white_presentation(output_path="NEXORA_Clean_White_Presentation.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Executive Light / White Theme Palette
    COLOR_BG = RGBColor(248, 250, 252)          # #f8fafc (Clean Slate Light Background)
    COLOR_CARD = RGBColor(255, 255, 255)        # #ffffff (Pure White Card)
    COLOR_CARD_BORDER = RGBColor(226, 232, 240)    # #e2e8f0 (Soft Border)
    
    # Modern Professional Accent Colors
    COLOR_PRIMARY = RGBColor(15, 23, 42)        # #0f172a (Deep Slate / Dark Navy)
    COLOR_BODY = RGBColor(51, 65, 85)           # #334155 (Slate 700)
    COLOR_MUTED = RGBColor(100, 116, 139)       # #64748b (Slate 500)
    
    COLOR_BRAND = RGBColor(79, 70, 229)         # #4f46e5 (Indigo 600)
    COLOR_CYAN = RGBColor(2, 132, 199)          # #0284c7 (Sky/Blue)
    COLOR_EMERALD = RGBColor(13, 148, 136)      # #0d9488 (Teal/Emerald)
    COLOR_AMBER = RGBColor(217, 119, 6)         # #d97706 (Amber)

    def add_blank_slide():
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = COLOR_BG
        bg_shape.line.color.rgb = COLOR_BG
        return slide

    def add_header(slide, tag_text, title_text, subtitle_text=None):
        tx_box = slide.shapes.add_textbox(Inches(0.9), Inches(0.45), Inches(11.5), Inches(1.2))
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p_tag = tf.paragraphs[0]
        p_tag.text = tag_text.upper()
        p_tag.font.size = Pt(10)
        p_tag.font.bold = True
        p_tag.font.color.rgb = COLOR_BRAND
        p_tag.space_after = Pt(3)

        p_title = tf.add_paragraph()
        p_title.text = title_text
        p_title.font.size = Pt(22)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_PRIMARY
        
        if subtitle_text:
            p_sub = tf.add_paragraph()
            p_sub.text = subtitle_text
            p_sub.font.size = Pt(11.5)
            p_sub.font.color.rgb = COLOR_MUTED
            p_sub.space_before = Pt(2)

    # ----------------------------------------------------
    # SLIDE 1: Title Slide
    # ----------------------------------------------------
    s1 = add_blank_slide()

    top_bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.1))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLOR_BRAND
    top_bar.line.color.rgb = COLOR_BRAND

    hero_card = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.0), Inches(11.333), Inches(5.6))
    hero_card.fill.solid()
    hero_card.fill.fore_color.rgb = COLOR_CARD
    hero_card.line.color.rgb = COLOR_CARD_BORDER

    tf1 = hero_card.text_frame
    tf1.word_wrap = True
    tf1.margin_left = Inches(0.6)
    tf1.margin_right = Inches(0.6)
    tf1.margin_top = Inches(0.5)

    p_badge = tf1.paragraphs[0]
    p_badge.text = "AI-POWERED EDUCATIONAL & CAREER PLATFORM"
    p_badge.font.size = Pt(11)
    p_badge.font.bold = True
    p_badge.font.color.rgb = COLOR_BRAND
    p_badge.space_after = Pt(10)

    p_h1 = tf1.add_paragraph()
    p_h1.text = "NEXORA SaaS"
    p_h1.font.size = Pt(44)
    p_h1.font.bold = True
    p_h1.font.color.rgb = COLOR_PRIMARY
    p_h1.space_after = Pt(6)

    p_sub = tf1.add_paragraph()
    p_sub.text = "AI-Powered Course Recommendation & Dynamic Roadmap Engine"
    p_sub.font.size = Pt(20)
    p_sub.font.bold = True
    p_sub.font.color.rgb = COLOR_CYAN
    p_sub.space_after = Pt(14)

    p_desc = tf1.add_paragraph()
    p_desc.text = "Automated Career Navigation, Topological Curriculum Sequencing, Diagnostic Skill Gap Scoring & 24/7 Context-Aware AI Guidance."
    p_desc.font.size = Pt(13)
    p_desc.font.color.rgb = COLOR_BODY
    p_desc.space_after = Pt(22)

    p_pres = tf1.add_paragraph()
    p_pres.text = "Presenter Information: Team Talent Innovators  |  Category: AI in Education & Career Placement"
    p_pres.font.size = Pt(12)
    p_pres.font.bold = True
    p_pres.font.color.rgb = COLOR_PRIMARY
    p_pres.space_after = Pt(6)

    p_meta = tf1.add_paragraph()
    p_meta.text = "Live Deployment: https://nexora-path-finder.vercel.app  •  Powered by Google Gemini AI & React 19"
    p_meta.font.size = Pt(10.5)
    p_meta.font.color.rgb = COLOR_EMERALD

    # ----------------------------------------------------
    # SLIDE 2: Presenter Details
    # ----------------------------------------------------
    s2 = add_blank_slide()
    add_header(s2, "Presenter Details", "Project Presentation Team", "The innovators behind NEXORA's design and execution")

    card_team = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.9), Inches(5.6), Inches(4.9))
    card_team.fill.solid()
    card_team.fill.fore_color.rgb = COLOR_CARD
    card_team.line.color.rgb = COLOR_CARD_BORDER

    tf_team = card_team.text_frame
    tf_team.word_wrap = True
    tf_team.margin_left = Inches(0.4)
    tf_team.margin_right = Inches(0.4)
    tf_team.margin_top = Inches(0.4)

    p = tf_team.paragraphs[0]
    p.text = "TEAM DETAILS"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_BRAND
    p.space_after = Pt(8)

    p = tf_team.add_paragraph()
    p.text = "Team Name:"
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_MUTED
    
    p = tf_team.add_paragraph()
    p.text = "Talent Innovators"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(20)

    p = tf_team.add_paragraph()
    p.text = "Project Track / Theme:"
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_MUTED

    p = tf_team.add_paragraph()
    p.text = "AI in Education, Adaptive Learning & Career Placement"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

    card_members = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.9), Inches(5.65), Inches(4.9))
    card_members.fill.solid()
    card_members.fill.fore_color.rgb = COLOR_CARD
    card_members.line.color.rgb = COLOR_CARD_BORDER

    tf_mem = card_members.text_frame
    tf_mem.word_wrap = True
    tf_mem.margin_left = Inches(0.4)
    tf_mem.margin_right = Inches(0.4)
    tf_mem.margin_top = Inches(0.4)

    p = tf_mem.paragraphs[0]
    p.text = "TEAM MEMBERS"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_EMERALD
    p.space_after = Pt(14)

    members_placeholder = [
        "1.  ___________________________________  (Team Lead)",
        "2.  ___________________________________  (Developer / AI)",
        "3.  ___________________________________  (Developer / UI)",
        "4.  ___________________________________  (Research / QA)"
    ]

    for m in members_placeholder:
        p = tf_mem.add_paragraph()
        p.text = m
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_PRIMARY
        p.space_after = Pt(16)

    # ----------------------------------------------------
    # SLIDE 3: Problem Understanding
    # ----------------------------------------------------
    s3 = add_blank_slide()
    add_header(s3, "Problem Understanding", "The Gap in Modern Technical Learning & Career Readiness", "Understanding the structural pain points that cause 90% of learners to drop out")

    cards_s3 = [
        ("The Gap in E-Learning", 
         "Online learners face severe information overload with thousands of disconnected video tutorials. Finding the exact right courses and sequencing them into an actionable study path is frustrating, time-consuming, and prone to abandonment.",
         COLOR_BRAND),
        ("Core Pain Points", 
         "• Lack of Goal Alignment: Generic courses ignore student starting baselines.\n• Absence of Prerequisite Clarity: Students fail advanced topics due to unassessed foundation gaps.\n• No Accountability: Zero real-time diagnostic feedback or personalized next steps.",
         COLOR_CYAN),
        ("Our Core Mission", 
         "Build an intelligent SaaS platform that dynamically recommends courses, explains technical relevance using Google Gemini LLMs, visualizes prerequisite dependency DAGs, and generates adaptive daily milestones.",
         COLOR_EMERALD)
    ]

    for i, (title, desc, col) in enumerate(cards_s3):
        left = Inches(0.9 + i * 3.9)
        top = Inches(1.8)
        
        box = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.7), Inches(5.1))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_CARD
        box.line.color.rgb = COLOR_CARD_BORDER

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.3)
        tf.margin_right = Inches(0.3)
        tf.margin_top = Inches(0.35)

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = col
        p.space_after = Pt(12)

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_BODY
        p2.line_spacing = 1.25

    # ----------------------------------------------------
    # SLIDE 4: Solution Approach
    # ----------------------------------------------------
    s4 = add_blank_slide()
    add_header(s4, "Solution Approach", "How NEXORA Solves Curriculum Sequencing", "An end-to-end adaptive framework guiding learners from goal intake to recruiter readiness")

    solutions = [
        ("1. Smart Goal Matching", 
         "Synthesizes natural language inputs (voice/text) into tailored goal profiles (SWE, AI/ML, JEE, NEET, Data Science, Career Switch), matching learners directly to curated high-yield topics.",
         COLOR_BRAND),
        ("2. Real-Time Explanation Engine", 
         "Curriculum analysis breaking down every milestone into prerequisites, core conceptual topics, verified free learning resources, and capstone project blueprints.",
         COLOR_CYAN),
        ("3. Dynamic DAG Sequencing", 
         "Generates topological Directed Acyclic Graphs with strict node status tracking (Completed, In Progress, Unlocked, Locked, Remediation) based on real-time diagnostic performance.",
         COLOR_EMERALD),
        ("4. Secure Workspaces & ATS Readiness", 
         "Zero-latency client persistent session with optional Firebase Cloud sync, paired with an automated ATS resume builder calculating real-time recruiter match scores.",
         COLOR_PRIMARY)
    ]

    for i, (title, desc, col) in enumerate(solutions):
        col_idx = i % 2
        row_idx = i // 2
        left = Inches(0.9 + col_idx * 5.85)
        top = Inches(1.8 + row_idx * 2.55)
        
        box = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.65), Inches(2.35))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_CARD
        box.line.color.rgb = COLOR_CARD_BORDER

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.3)
        tf.margin_right = Inches(0.3)
        tf.margin_top = Inches(0.25)

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14.5)
        p.font.bold = True
        p.font.color.rgb = col
        p.space_after = Pt(8)

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_BODY
        p2.line_spacing = 1.2

    # ----------------------------------------------------
    # SLIDE 5: Visual Architecture Diagram (Dedicated Architecture Slide)
    # ----------------------------------------------------
    s5 = add_blank_slide()
    add_header(s5, "System Architecture Diagram", "Visual End-to-End System Architecture", "Dataflow across UI presentation, adaptive engines, Gemini LLM reasoning, and cloud persistence")

    diag_path = "public/nexora_architecture_diagram.png"
    if os.path.exists(diag_path):
        s5.shapes.add_picture(diag_path, Inches(0.9), Inches(1.7), Inches(11.533), Inches(5.3))
    else:
        # Fallback if image not found
        arch_box = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.7), Inches(11.533), Inches(5.3))
        arch_box.fill.solid()
        arch_box.fill.fore_color.rgb = COLOR_CARD
        arch_box.line.color.rgb = COLOR_CARD_BORDER

    # ----------------------------------------------------
    # SLIDE 6: System Architecture Layer Breakdown
    # ----------------------------------------------------
    s6 = add_blank_slide()
    add_header(s6, "Technical Execution", "Architectural Stack & System Components", "Detailed technical specifications of the 4 core layers")

    arch_layers = [
        ("Frontend Application Layer", 
         "• React 19 & TypeScript for strict end-to-end type safety.\n• TailwindCSS with clean executive design tokens and Lucide Icons.\n• Web Speech API for voice requirement gathering & chat speech-to-text.\n• Responsive single-page application with zero-latency tab navigation.",
         COLOR_BRAND),
        ("Adaptive Engine & Algorithms", 
         "• Directed Acyclic Graph (DAG) Engine with topological prerequisite sorting.\n• Automated Diagnostic Skill Gap Scorer calculating mastery vs benchmarks.\n• Next-Best-Action (NBA) Scorer computing single highest-leverage task.\n• What-If Monte-Carlo Timeline Simulator modeling pacing shifts.",
         COLOR_CYAN),
        ("AI Services & LLM Layer", 
         "• Google Gemini 1.5 Flash & 2.0 API with multi-turn conversational reasoning.\n• Dynamic system prompt pipeline injected with student profile & milestones.\n• Resilient offline semantic intent classifier for instant fallback answers.\n• Dynamic MCQ assessment generator for custom skills.",
         COLOR_EMERALD),
        ("Storage, Export & Deployment", 
         "• LocalStorage Persistent Engine for instant offline session caching.\n• Firebase Cloud Auth & Cloud Firestore profile sync (Optional).\n• jsPDF & jsPDF-AutoTable for instantaneous ATS Resume PDF generation.\n• Continuous Deployment Pipeline: GitHub auto-deployed via Vercel.",
         COLOR_PRIMARY)
    ]

    for i, (title, desc, col) in enumerate(arch_layers):
        col_idx = i % 2
        row_idx = i // 2
        left = Inches(0.9 + col_idx * 5.85)
        top = Inches(1.8 + row_idx * 2.55)
        
        box = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.65), Inches(2.35))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_CARD
        box.line.color.rgb = COLOR_CARD_BORDER

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.3)
        tf.margin_right = Inches(0.3)
        tf.margin_top = Inches(0.22)

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = col
        p.space_after = Pt(6)

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = COLOR_BODY
        p2.line_spacing = 1.18

    # ----------------------------------------------------
    # SLIDE 7: AI & ML Techniques Used
    # ----------------------------------------------------
    s7 = add_blank_slide()
    add_header(s7, "AI & ML Techniques", "Core Machine Learning & LLM Implementations", "How generative intelligence and topological graphs deliver customized learning")

    ai_techniques = [
        ("1. Large Language Model Integration", 
         "Integrated Google Gemini 1.5 Flash and 2.0 API with dynamic context injection. Sends full user state (Target Goal, Completed Milestones, Active Skill Gaps, Pacing) to generate tailored conversational guidance, code examples, and doubt resolution.",
         COLOR_BRAND),
        ("2. Topological Graph Traversal (DAG)", 
         "Implements graph-theoretic dependency trees where skills represent nodes and prerequisites represent directed edges. Automatically triggers remedial sub-branches when assessment scores drop below threshold benchmarks (70%).",
         COLOR_CYAN),
        ("3. Structured Output & JSON Enforcement", 
         "Utilizes strict JSON schema prompting to dynamically generate 8-step roadmap milestones, goal-specific challenge vectors, and diagnostic multiple-choice assessments on the fly.",
         COLOR_EMERALD),
        ("4. Latency Management & Offline Intelligence", 
         "Features multi-tiered execution with fast 4-second timeout handling and a semantic intent classifier capable of addressing technical comparisons (e.g. Web Dev vs ML), OOP, and DSA offline without freezing.",
         COLOR_PRIMARY)
    ]

    for i, (title, desc, col) in enumerate(ai_techniques):
        col_idx = i % 2
        row_idx = i // 2
        left = Inches(0.9 + col_idx * 5.85)
        top = Inches(1.8 + row_idx * 2.55)
        
        box = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.65), Inches(2.35))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_CARD
        box.line.color.rgb = COLOR_CARD_BORDER

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.3)
        tf.margin_right = Inches(0.3)
        tf.margin_top = Inches(0.22)

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = col
        p.space_after = Pt(6)

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = COLOR_BODY
        p2.line_spacing = 1.2

    # ----------------------------------------------------
    # SLIDE 8: Key Features & Core Workflows
    # ----------------------------------------------------
    s8 = add_blank_slide()
    add_header(s8, "Product Features", "Key Features & Core User Workflows", "Comprehensive suite of integrated tools designed for continuous student mastery")

    features_s8 = [
        ("AI Conversational Onboarding", "Interactive voice & text interface gathering user goals, time availability, and background.", COLOR_BRAND),
        ("Interactive Flowchart DAG", "Visual dependency roadmap with clickable nodes, learning objectives, and status tags.", COLOR_CYAN),
        ("Automated Diagnostic Engine", "Benchmark-driven MCQ quizzes that calculate real-time mastery percentages per skill.", COLOR_EMERALD),
        ("Next-Best-Action (NBA)", "Algorithmic decision engine computing the single highest-yield study task for today.", COLOR_PRIMARY),
        ("24/7 Context-Aware AI Chatbot", "Intelligent tutor providing code blocks with copy buttons, comparisons, and exam tips.", COLOR_BRAND),
        ("ATS Resume Builder & Scoring", "Instant ATS score calculation, Google XYZ bullet optimizers, and 1-click PDF download.", COLOR_CYAN),
        ("What-If Scenario Simulator", "Interactive time-commitment slider modeling how study pacing impacts target completion.", COLOR_EMERALD),
        ("100% Free Resources Catalog", "Curated zero-cost materials from MIT OpenCourseWare, CS50, Disha, HCV, & freeCodeCamp.", COLOR_PRIMARY)
    ]

    for i, (title, desc, col) in enumerate(features_s8):
        col_idx = i % 4
        row_idx = i // 4
        left = Inches(0.9 + col_idx * 2.9)
        top = Inches(1.8 + row_idx * 2.55)
        
        box = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.75), Inches(2.35))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_CARD
        box.line.color.rgb = COLOR_CARD_BORDER

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.2)

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = col
        p.space_after = Pt(6)

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(10)
        p2.font.color.rgb = COLOR_BODY
        p2.line_spacing = 1.15

    # ----------------------------------------------------
    # SLIDE 9: Summary & Conclusion
    # ----------------------------------------------------
    s9 = add_blank_slide()

    card_s9 = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.0), Inches(11.333), Inches(5.6))
    card_s9.fill.solid()
    card_s9.fill.fore_color.rgb = COLOR_CARD
    card_s9.line.color.rgb = COLOR_CARD_BORDER

    tf9 = card_s9.text_frame
    tf9.word_wrap = True
    tf9.margin_left = Inches(0.6)
    tf9.margin_right = Inches(0.6)
    tf9.margin_top = Inches(0.5)

    p = tf9.paragraphs[0]
    p.text = "THE FUTURE OF ADAPTIVE EDUCATION"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_BRAND
    p.space_after = Pt(8)

    p1 = tf9.add_paragraph()
    p1.text = "Thank You! Experience NEXORA SaaS Live"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_PRIMARY
    p1.space_after = Pt(12)

    p2 = tf9.add_paragraph()
    p2.text = "NEXORA bridges the gap between disorganized online educational content and structured, recruiter-ready student outcomes through intelligent sequencing, real-time diagnostics, and personalized AI mentorship."
    p2.font.size = Pt(13)
    p2.font.color.rgb = COLOR_BODY
    p2.space_after = Pt(24)

    p3 = tf9.add_paragraph()
    p3.text = "• Team: Talent Innovators\n• Live Web Application: https://nexora-path-finder.vercel.app\n• GitHub Repository: https://github.com/Sirivennela310505/NEXORA"
    p3.font.size = Pt(12)
    p3.font.bold = True
    p3.font.color.rgb = COLOR_BRAND

    prs.save(output_path)
    print(f"White presentation with visual architecture successfully saved to: {output_path}")

if __name__ == "__main__":
    create_nexora_white_presentation()
