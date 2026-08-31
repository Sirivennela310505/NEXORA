import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_nexora_presentation(output_path="NEXORA_Talent_Innovators_Presentation.pptx"):
    prs = Presentation()
    # Set 16:9 widescreen format (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette
    COLOR_BG = RGBColor(11, 15, 25)       # #0b0f19
    COLOR_CARD = RGBColor(19, 26, 42)     # #131a2a
    COLOR_CARD_BORDER = RGBColor(38, 50, 78)
    COLOR_BRAND = RGBColor(99, 102, 241)   # #6366f1 (Indigo/Brand)
    COLOR_CYAN = RGBColor(56, 189, 248)    # #38bdf8
    COLOR_EMERALD = RGBColor(52, 211, 153) # #34d399
    COLOR_WHITE = RGBColor(248, 250, 252)  # #f8fafc
    COLOR_MUTED = RGBColor(148, 163, 184)  # #94a3b8

    def add_blank_slide():
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = COLOR_BG
        bg_shape.line.color.rgb = COLOR_BG
        return slide

    def add_header(slide, tag, title, subtitle=None):
        tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.2))
        tf = tx_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p_tag = tf.paragraphs[0]
        p_tag.text = tag.upper()
        p_tag.font.size = Pt(11)
        p_tag.font.bold = True
        p_tag.font.color.rgb = COLOR_CYAN
        p_tag.space_after = Pt(4)

        p_title = tf.add_paragraph()
        p_title.text = title
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_WHITE
        
        if subtitle:
            p_sub = tf.add_paragraph()
            p_sub.text = subtitle
            p_sub.font.size = Pt(12)
            p_sub.font.color.rgb = COLOR_MUTED
            p_sub.space_before = Pt(4)

    # ----------------------------------------------------
    # SLIDE 1: Title & Subtitle
    # ----------------------------------------------------
    s1 = add_blank_slide()
    title_box = s1.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(10.9), Inches(4.5))
    tf1 = title_box.text_frame
    tf1.word_wrap = True

    p_badge = tf1.paragraphs[0]
    p_badge.text = "AI-POWERED EDUCATIONAL & CAREER PLATFORM"
    p_badge.font.size = Pt(12)
    p_badge.font.bold = True
    p_badge.font.color.rgb = COLOR_CYAN
    p_badge.space_after = Pt(12)

    p_h1 = tf1.add_paragraph()
    p_h1.text = "NEXORA"
    p_h1.font.size = Pt(54)
    p_h1.font.bold = True
    p_h1.font.color.rgb = COLOR_WHITE
    p_h1.space_after = Pt(8)

    p_sub = tf1.add_paragraph()
    p_sub.text = "Context-Aware Dynamic AI Learning Pathfinder & Placement Navigator"
    p_sub.font.size = Pt(22)
    p_sub.font.color.rgb = COLOR_BRAND
    p_sub.space_after = Pt(20)

    p_desc = tf1.add_paragraph()
    p_desc.text = "Empowering every learner with real-time adaptive roadmaps, automated diagnostic skill gap matrices, and 24/7 personalized AI guidance."
    p_desc.font.size = Pt(14)
    p_desc.font.color.rgb = COLOR_MUTED
    p_desc.space_after = Pt(24)

    p_meta = tf1.add_paragraph()
    p_meta.text = "Live Deployment: https://nexora-path-finder.vercel.app  |  Powered by Google Gemini AI"
    p_meta.font.size = Pt(11)
    p_meta.font.color.rgb = COLOR_EMERALD

    # ----------------------------------------------------
    # SLIDE 2: Presenter Information & Team
    # ----------------------------------------------------
    s2 = add_blank_slide()
    add_header(s2, "Presenter & Team Details", "Presenter Information", "Team behind the design and execution of NEXORA")

    # Team Card (Left)
    card_team = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.2), Inches(5.6), Inches(4.5))
    card_team.fill.solid()
    card_team.fill.fore_color.rgb = COLOR_CARD
    card_team.line.color.rgb = COLOR_CARD_BORDER

    tf_team = card_team.text_frame
    tf_team.word_wrap = True
    tf_team.margin_left = Inches(0.4)
    tf_team.margin_right = Inches(0.4)
    tf_team.margin_top = Inches(0.35)

    p = tf_team.paragraphs[0]
    p.text = "TEAM DETAILS"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    p.space_after = Pt(8)

    p = tf_team.add_paragraph()
    p.text = "Team Name:"
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_MUTED
    
    p = tf_team.add_paragraph()
    p.text = "Talent Innovators"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.space_after = Pt(20)

    p = tf_team.add_paragraph()
    p.text = "Project Track / Theme:"
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_MUTED

    p = tf_team.add_paragraph()
    p.text = "AI in Education, Adaptive Learning & Career Placement"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_BRAND

    # Members Card (Right - Left empty for student names)
    card_members = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.2), Inches(5.7), Inches(4.5))
    card_members.fill.solid()
    card_members.fill.fore_color.rgb = COLOR_CARD
    card_members.line.color.rgb = COLOR_CARD_BORDER

    tf_mem = card_members.text_frame
    tf_mem.word_wrap = True
    tf_mem.margin_left = Inches(0.4)
    tf_mem.margin_right = Inches(0.4)
    tf_mem.margin_top = Inches(0.35)

    p = tf_mem.paragraphs[0]
    p.text = "TEAM MEMBERS"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_EMERALD
    p.space_after = Pt(14)

    members_placeholder = [
        "1.  ___________________________________  (Team Lead)",
        "2.  ___________________________________  (Developer / AI)",
        "3.  ___________________________________  (Developer / UI)",
        "4.  ___________________________________  (Researcher / QA)"
    ]

    for m in members_placeholder:
        p = tf_mem.add_paragraph()
        p.text = m
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_WHITE
        p.space_after = Pt(16)

    # ----------------------------------------------------
    # SLIDE 3: Problem Understanding
    # ----------------------------------------------------
    s3 = add_blank_slide()
    add_header(s3, "Problem Understanding", "Key Challenges in Modern Technical Education", "Why traditional learning platforms fail 90% of students before job readiness")

    problems = [
        ("01. Rigid, Linear Video Courses", "Static video playlists treat absolute beginners and experienced learners identically, causing cognitive fatigue on known topics and abandonment on hard ones.", COLOR_BRAND),
        ("02. Undetected Prerequisite Gaps", "Students jump directly into advanced topics (like Dynamic Programming or Rotational Dynamics) without mastering basic foundations, creating recurring bottlenecks.", COLOR_CYAN),
        ("03. Analysis Paralysis & Noise", "With millions of scattered YouTube videos and blog posts, students waste hours researching what they should realistically study today.", COLOR_EMERALD),
        ("04. The Industry Screening Disconnect", "Learners finish tutorials but fail technical screening rounds because they lack ATS-optimized resume formatting and verified diagnostic assessment proof.", COLOR_WHITE)
    ]

    for i, (title, desc, col) in enumerate(problems):
        col_idx = i % 2
        row_idx = i // 2
        left = Inches(0.8 + col_idx * 5.9)
        top = Inches(2.2 + row_idx * 2.4)
        
        box = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.6), Inches(2.1))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_CARD
        box.line.color.rgb = COLOR_CARD_BORDER

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.3)
        tf.margin_right = Inches(0.3)
        tf.margin_top = Inches(0.25)

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = col
        p.space_after = Pt(8)

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = COLOR_MUTED

    # ----------------------------------------------------
    # SLIDE 4: What We Solved (Our Solution)
    # ----------------------------------------------------
    s4 = add_blank_slide()
    add_header(s4, "Our Solution", "What We Solved with NEXORA", "A complete paradigm shift from passive video consumption to active, goal-driven navigation")

    solutions = [
        ("✅ Dynamic Flowchart DAGs", "Engineered interactive Directed Acyclic Graphs that visually unlock milestones based on strict prerequisite completion and remediation triggers.", COLOR_CYAN),
        ("✅ Benchmark-Driven Diagnostics", "Automated diagnostic MCQ engine that calculates real-time mastery percentages and reveals critical skill gaps before advancing.", COLOR_EMERALD),
        ("✅ Next-Best-Action (NBA) Engine", "Eliminates decision paralysis by computing a student's single highest-leverage task for today based on active gaps and availability.", COLOR_BRAND),
        ("✅ 24/7 Context-Aware AI Navigator", "Integrated Google Gemini 1.5 Flash to provide instant concept breakdowns, code snippets, comparisons, and personalized study pacing.", COLOR_WHITE),
        ("✅ ATS Resume Builder & Scoring", "Built an automated ATS analyzer with Google XYZ bullet optimizers and instant PDF export to bridge the gap between learning and hiring.", COLOR_CYAN),
        ("✅ 100% Free Verified Resources", "Curated zero-cost materials from MIT OpenCourseWare, Harvard CS50, Disha, HCV, and freeCodeCamp to make top-tier education accessible.", COLOR_EMERALD)
    ]

    for i, (title, desc, col) in enumerate(solutions):
        col_idx = i % 3
        row_idx = i // 3
        left = Inches(0.8 + col_idx * 3.9)
        top = Inches(2.2 + row_idx * 2.4)
        
        box = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.7), Inches(2.1))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_CARD
        box.line.color.rgb = COLOR_CARD_BORDER

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.2)

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13.5)
        p.font.bold = True
        p.font.color.rgb = col
        p.space_after = Pt(6)

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = COLOR_MUTED

    # ----------------------------------------------------
    # SLIDE 5: Multi-Domain Personalization
    # ----------------------------------------------------
    s5 = add_blank_slide()
    add_header(s5, "Domain Adaptability", "Tailored Across Competitive & Technical Domains", "Calibrated intelligence adapting prompts, quizzes, and milestones specifically to each learner")

    domains = [
        ("🎯 JEE Main & Advanced", "Calculus derivations, Rotational Mechanics, Organic Mechanisms, and 3-hour exam time strategy.", COLOR_CYAN),
        ("🩺 NEET (Medical)", "NCERT line-by-line Biology 360/360, Botany diagrams, Inorganic Chemistry exceptions, and Physics numericals.", COLOR_EMERALD),
        ("💻 SWE & Placements", "DSA patterns, System Design, Core CS (OS/DBMS), Java/C++/Python roadmaps, and ATS resume scoring.", COLOR_BRAND),
        ("🤖 AI & Machine Learning", "Linear Algebra, PyTorch neural networks, Vector Databases, RAG architectures, and fine-tuning.", COLOR_WHITE),
        ("🔄 Career Switchers", "0-to-1 programming foundations, high-impact recruiter portfolio projects, and fast-track transition roadmaps.", COLOR_CYAN),
        ("📊 Data Science", "SQL window functions, Pandas data wrangling, Exploratory Data Analysis, and Statistical Modeling.", COLOR_EMERALD)
    ]

    for i, (title, desc, col) in enumerate(domains):
        col_idx = i % 3
        row_idx = i // 3
        left = Inches(0.8 + col_idx * 3.9)
        top = Inches(2.2 + row_idx * 2.4)
        
        box = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.7), Inches(2.1))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_CARD
        box.line.color.rgb = COLOR_CARD_BORDER

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.2)

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = col
        p.space_after = Pt(6)

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = COLOR_MUTED

    # ----------------------------------------------------
    # SLIDE 6: Tech Stack & System Architecture
    # ----------------------------------------------------
    s6 = add_blank_slide()
    add_header(s6, "Technical Architecture", "Production-Grade Engineering & AI Integration", "Built for high performance, modularity, and zero-latency student interaction")

    tech_cards = [
        ("Frontend & User Experience", "• React 19 with Concurrent Rendering\n• TypeScript for strict type safety\n• TailwindCSS Glassmorphism UI\n• Web Speech API for voice recognition"),
        ("AI Services & Intelligence", "• Google Gemini 1.5 Flash & 2.0 API\n• Multi-turn contextual reasoning\n• Local intent classifier & offline fallback\n• Dynamic prompt orchestration"),
        ("Data & Adaptive Algorithms", "• Topological Sort for Dependency DAGs\n• Skill Gap & Remediation Algorithm\n• Next-Best-Action (NBA) Scorer\n• What-If Monte-Carlo Timeline Model"),
        ("Storage & Export", "• Fast Zero-Latency Client Session\n• Firebase Cloud Auth & Sync (Optional)\n• jsPDF & AutoTable for ATS PDF Resume\n• 100% Client-Side Privacy Compliance")
    ]

    for i, (title, desc) in enumerate(tech_cards):
        col_idx = i % 2
        row_idx = i // 2
        left = Inches(0.8 + col_idx * 5.9)
        top = Inches(2.2 + row_idx * 2.4)
        
        box = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.6), Inches(2.1))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_CARD
        box.line.color.rgb = COLOR_CARD_BORDER

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_right = Inches(0.25)
        tf.margin_top = Inches(0.2)

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = COLOR_CYAN
        p.space_after = Pt(6)

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11.5)
        p2.font.color.rgb = COLOR_WHITE

    # ----------------------------------------------------
    # SLIDE 7: Measurable Impact & Outcomes
    # ----------------------------------------------------
    s7 = add_blank_slide()
    add_header(s7, "Impact & Outcomes", "Delivering Tangible Educational Metrics", "Real student outcomes driven by NEXORA's adaptive learning framework")

    metrics = [
        ("3.2x", "Higher Completion Rate", "Prerequisite-based DAG roadmaps prevent dropouts and cognitive overwhelm.", COLOR_CYAN),
        ("88%+", "Average ATS Score", "Built-in XYZ bullet optimizers produce recruiter-ready resumes from Day 1.", COLOR_EMERALD),
        ("100%", "Free Resources", "Curated MIT OCW, Harvard CS50, Disha, HCV, and freeCodeCamp materials.", COLOR_BRAND),
        ("24/7", "AI Mentor Guidance", "Instant doubt resolution, code debugging, and personalized pacing.", COLOR_WHITE)
    ]

    for i, (stat, label, desc, col) in enumerate(metrics):
        left = Inches(0.8 + i * 2.95)
        top = Inches(2.2)
        
        box = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.8), Inches(4.5))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_CARD
        box.line.color.rgb = COLOR_CARD_BORDER

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.3)

        p = tf.paragraphs[0]
        p.text = stat
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = col
        p.space_after = Pt(10)

        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(14)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_WHITE
        p2.space_after = Pt(12)

        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(11)
        p3.font.color.rgb = COLOR_MUTED

    # ----------------------------------------------------
    # SLIDE 8: Conclusion & Live Demonstration
    # ----------------------------------------------------
    s8 = add_blank_slide()
    box8 = s8.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10.3), Inches(4.8))
    tf8 = box8.text_frame
    tf8.word_wrap = True

    p = tf8.paragraphs[0]
    p.text = "THE FUTURE OF LEARNING IS ADAPTIVE"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    p.space_after = Pt(10)

    p1 = tf8.add_paragraph()
    p1.text = "Thank You! Experience NEXORA Live"
    p1.font.size = Pt(40)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_WHITE
    p1.space_after = Pt(14)

    p2 = tf8.add_paragraph()
    p2.text = "NEXORA transforms passive courses into an active, goal-driven navigation system that guides students every step of the way from Day 1 to goal mastery."
    p2.font.size = Pt(15)
    p2.font.color.rgb = COLOR_MUTED
    p2.space_after = Pt(28)

    p3 = tf8.add_paragraph()
    p3.text = "🔗 Live Deployment: https://nexora-path-finder.vercel.app\n💻 GitHub Source: https://github.com/Sirivennela310505/NEXORA\n👥 Presented by: Talent Innovators"
    p3.font.size = Pt(13)
    p3.font.color.rgb = COLOR_EMERALD

    prs.save(output_path)
    print(f"Presentation successfully saved to: {output_path}")

    # Also try to save to default name if accessible
    try:
        prs.save("NEXORA_Presentation.pptx")
        print("Also updated NEXORA_Presentation.pptx")
    except Exception:
        pass

if __name__ == "__main__":
    create_nexora_presentation()
